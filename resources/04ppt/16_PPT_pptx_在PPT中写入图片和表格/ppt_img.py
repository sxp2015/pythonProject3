# -*- coding: utf-8 -*-
# ⭐Python 60套学习资源：http://t.cn/A6xASARf
# 📱公众号 :程序员晚枫 读者交流群：http://www.python4office.cn/wechat-group/
# 🏠2022最新视频：1行代码，实现自动化办公：https://www.bilibili.com/video/BV1pT4y1k7FH


# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches,Pt

ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[1])# 在PPT中插入一个幻灯片

# left = Inches(1)
# top = Inches(1)
# width = Inches(2)
# height = Inches(2)
#
# img = slide.shapes.add_picture('img.jpg',left,top,width,height)

rows = 2
cols = 2
left = Inches(1)
top = Inches(1)
width = Inches(4)
height = Inches(4)

table = slide.shapes.add_table(rows,cols,left,top,width,height).table
table.columns[0].width = Inches(1)
table.columns[1].width = Inches(3)
table.cell(0,0).text = '00'
table.cell(0,1).text = '01'
table.cell(1,0).text = '10'
table.cell(1,1).text = '11' #二进制的11，代表十进制的多少？

ppt.save('text.pptx')
