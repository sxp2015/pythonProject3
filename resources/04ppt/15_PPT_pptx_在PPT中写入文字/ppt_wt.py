# -*- coding: utf-8 -*-
# ⭐Python 60套学习资源：http://t.cn/A6xASARf
# 📱公众号 :程序员晚枫 读者交流群：http://www.python4office.cn/wechat-group/
# 🏠2022最新视频：1行代码，实现自动化办公：https://www.bilibili.com/video/BV1pT4y1k7FH


# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches,Pt

ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[1])# 在PPT中插入一个幻灯片

body_shape = slide.shapes.placeholders
# body_shape[0].text = '这是占位符0'
# body_shape[1].text = '这是占位符1'
#
title_shape = slide.shapes.title
title_shape.text = '这是标题'
# subtitle = slide.shapes.placeholders[1] #取出本页第二个文本框
# subtitle.text = '这是文本框'
#
# new_paragraph = body_shape[1].text_frame.add_paragraph()
# new_paragraph.text = '新段落'
# new_paragraph.font.bold = True
# new_paragraph.font.italic = True
# new_paragraph.font.size = Pt(15)
# new_paragraph.font.underline = True
#
left = Inches(2)
top = Inches(2)
width = Inches(3)
height = Inches(3)
#
#
#
textbox = slide.shapes.add_textbox(left,top,width,height)
textbox.text = 'new textbox'
# 如何在文本框里添加第二段文字？
# new_para = textbox.text_frame.add_paragraph()
# new_para.text = '第二段文字'
#


ppt.save('test.pptx')